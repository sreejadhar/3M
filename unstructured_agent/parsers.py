"""
Format parsers for the unstructured data intelligence agent.

Each parser extracts:
  1. Structural metadata (title, author, page count, etc.)
  2. A text window — the ephemeral excerpt used for LLM semantic extraction.

Raw text is NEVER persisted. It is returned in-memory only for the duration
of the extraction pipeline and then discarded.
"""
from __future__ import annotations

import io
import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional

logger = logging.getLogger(__name__)

# Token estimate: 1 token ≈ 4 chars (conservative)
_APPROX_TOKENS = lambda text: len(text) // 4

_WINDOW_LEAD_CHARS  = 8000   # ~2000 tokens from the start
_WINDOW_TRAIL_CHARS = 1200   # ~300 tokens from the end
_SHORT_DOC_CHARS    = 2000   # ~500 tokens — use full text below this


@dataclass
class ParseResult:
    file_type: str
    # Structural metadata
    title: str = ""
    author: str = ""
    page_count: Optional[int] = None
    word_count: Optional[int] = None
    section_headers: List[str] = field(default_factory=list)
    # Extraction quality signals
    ocr_used: bool = False
    ocr_confidence: Optional[float] = None
    # The ephemeral text window (never persisted)
    text_window: str = ""
    parse_error: Optional[str] = None


def _build_window(full_text: str) -> str:
    if len(full_text) <= _SHORT_DOC_CHARS:
        return full_text
    lead  = full_text[:_WINDOW_LEAD_CHARS]
    trail = full_text[-_WINDOW_TRAIL_CHARS:]
    return lead + "\n…\n" + trail


def parse_file(file_path: str) -> ParseResult:
    ext = Path(file_path).suffix.lower()
    try:
        if ext == ".pdf":
            return _parse_pdf(file_path)
        if ext in (".docx", ".doc"):
            return _parse_docx(file_path)
        if ext in (".pptx", ".ppt"):
            return _parse_pptx(file_path)
        if ext in (".xlsx", ".xls"):
            return _parse_xlsx(file_path)
        if ext in (".html", ".htm"):
            return _parse_html(file_path)
        if ext in (".md", ".rst"):
            return _parse_md(file_path)
        if ext in (".txt", ".csv"):
            return _parse_txt(file_path)
    except Exception as exc:
        logger.warning("Parse error for %s: %s", file_path, exc)
        return ParseResult(
            file_type=ext.lstrip(".") or "unknown",
            parse_error=str(exc),
        )
    return ParseResult(file_type="unknown", parse_error="Unsupported format")


# ── PDF ───────────────────────────────────────────────────────────────────────

def _parse_pdf(path: str) -> ParseResult:
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        return ParseResult(file_type="pdf", parse_error="pdfplumber not installed")

    result = ParseResult(file_type="pdf")
    texts: List[str] = []
    headers: List[str] = []

    with pdfplumber.open(path) as pdf:
        result.page_count = len(pdf.pages)
        if pdf.metadata:
            result.title  = (pdf.metadata.get("Title")  or "").strip()
            result.author = (pdf.metadata.get("Author") or "").strip()

        for page in pdf.pages:
            page_text = page.extract_text() or ""
            if page_text.strip():
                texts.append(page_text)

    full_text = "\n".join(texts)

    # If no text extracted from digital PDF, try OCR
    if not full_text.strip():
        result.ocr_used = True
        full_text, result.ocr_confidence = _ocr_pdf(path)

    result.word_count = len(full_text.split())
    result.section_headers = _extract_headers_heuristic(full_text)
    result.text_window = _build_window(full_text)
    return result


def _ocr_pdf(path: str):
    try:
        import pdf2image  # type: ignore
        import pytesseract  # type: ignore
    except ImportError:
        return "", None

    images = pdf2image.convert_from_path(path, dpi=150)
    texts = []
    confidences = []
    for img in images[:20]:  # cap at 20 pages for OCR to control cost
        data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
        texts.append(" ".join(w for w in data["text"] if w.strip()))
        confs = [c for c in data["conf"] if c != -1]
        if confs:
            confidences.extend(confs)
    avg_conf = (sum(confidences) / len(confidences)) if confidences else None
    return "\n".join(texts), (avg_conf / 100.0 if avg_conf is not None else None)


# ── DOCX ──────────────────────────────────────────────────────────────────────

def _parse_docx(path: str) -> ParseResult:
    try:
        import docx  # type: ignore  (python-docx)
    except ImportError:
        return ParseResult(file_type="docx", parse_error="python-docx not installed")

    doc = docx.Document(path)
    result = ParseResult(file_type="docx")

    props = doc.core_properties
    result.title  = (props.title  or "").strip()
    result.author = (props.author or "").strip()

    paragraphs: List[str] = []
    headers: List[str] = []
    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            headers.append(para.text.strip())
        if para.text.strip():
            paragraphs.append(para.text.strip())

    full_text = "\n".join(paragraphs)
    result.word_count = len(full_text.split())
    result.section_headers = headers
    result.text_window = _build_window(full_text)
    return result


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _parse_pptx(path: str) -> ParseResult:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return ParseResult(file_type="pptx", parse_error="python-pptx not installed")

    prs = Presentation(path)
    result = ParseResult(file_type="pptx")
    result.page_count = len(prs.slides)

    texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)

    full_text = "\n".join(texts)
    result.word_count = len(full_text.split())
    result.text_window = _build_window(full_text)
    return result


# ── XLSX ──────────────────────────────────────────────────────────────────────

def _parse_xlsx(path: str) -> ParseResult:
    try:
        import openpyxl  # type: ignore
    except ImportError:
        return ParseResult(file_type="xlsx", parse_error="openpyxl not installed")

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    result = ParseResult(file_type="xlsx")
    result.section_headers = wb.sheetnames

    # Use sheet names and named ranges as the text window (no cell values)
    parts = ["Sheets: " + ", ".join(wb.sheetnames)]
    if wb.defined_names:
        names = [dn.name for dn in wb.defined_names.definedName]
        if names:
            parts.append("Named ranges: " + ", ".join(names))
    result.text_window = "\n".join(parts)
    return result


# ── HTML ──────────────────────────────────────────────────────────────────────

def _parse_html(path: str) -> ParseResult:
    try:
        from bs4 import BeautifulSoup  # type: ignore
    except ImportError:
        return ParseResult(file_type="html", parse_error="beautifulsoup4 not installed")

    result = ParseResult(file_type="html")
    with open(path, encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    result.title = (soup.title.string or "").strip() if soup.title else ""

    headers = []
    for tag in soup.find_all(re.compile(r"^h[1-3]$")):
        h = tag.get_text(strip=True)
        if h:
            headers.append(h)
    result.section_headers = headers

    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    full_text = soup.get_text(separator="\n")
    full_text = re.sub(r"\n{3,}", "\n\n", full_text).strip()
    result.word_count = len(full_text.split())
    result.text_window = _build_window(full_text)
    return result


# ── Markdown / RST ────────────────────────────────────────────────────────────

def _parse_md(path: str) -> ParseResult:
    result = ParseResult(file_type="md")
    with open(path, encoding="utf-8", errors="replace") as f:
        text = f.read()

    headers = re.findall(r"^#{1,3}\s+(.+)$", text, re.MULTILINE)
    result.section_headers = [h.strip() for h in headers]
    result.word_count = len(text.split())
    result.text_window = _build_window(text)
    return result


# ── Plain text / CSV ─────────────────────────────────────────────────────────

def _parse_txt(path: str) -> ParseResult:
    result = ParseResult(file_type="txt")
    with open(path, encoding="utf-8", errors="replace") as f:
        text = f.read()
    result.word_count = len(text.split())
    result.text_window = _build_window(text)
    return result


# ── Utilities ─────────────────────────────────────────────────────────────────

def _extract_headers_heuristic(text: str, max_headers: int = 20) -> List[str]:
    """Rough header extraction from unstructured text using capitalisation signals."""
    candidates = []
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped or len(stripped) > 100:
            continue
        if stripped.isupper() and len(stripped) > 3:
            candidates.append(stripped.title())
        elif re.match(r"^\d+[\.\)]\s+[A-Z]", stripped):
            candidates.append(stripped)
        if len(candidates) >= max_headers:
            break
    return candidates
