"""
Plain-text extraction for uploaded documents.

Each format handler is best-effort: if the optional library it needs isn't
installed, it raises ExtractionError with a clear message instead of
crashing the pipeline.
"""
from __future__ import annotations

from pathlib import Path

_MAX_CHARS = 200_000  # cap stored/extracted text to keep the DB light


class ExtractionError(RuntimeError):
    pass


def _extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _extract_docx(path: Path) -> str:
    try:
        import docx
    except ImportError:
        raise ExtractionError("python-docx is not installed — run `pip install python-docx`")
    doc = docx.Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text for c in row.cells))
    return "\n".join(parts)


def _extract_pdf(path: Path) -> str:
    try:
        import pypdf
    except ImportError:
        raise ExtractionError("pypdf is not installed — run `pip install pypdf`")
    reader = pypdf.PdfReader(str(path))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_pptx(path: Path) -> str:
    try:
        import pptx
    except ImportError:
        raise ExtractionError("python-pptx is not installed — run `pip install python-pptx`")
    prs = pptx.Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        raise ExtractionError("openpyxl is not installed — run `pip install openpyxl`")
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"# {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


_HANDLERS = {
    ".txt": _extract_txt, ".md": _extract_txt, ".csv": _extract_txt,
    ".html": _extract_txt, ".htm": _extract_txt, ".rst": _extract_txt,
    ".docx": _extract_docx, ".doc": _extract_docx,
    ".pdf": _extract_pdf,
    ".pptx": _extract_pptx, ".ppt": _extract_pptx,
    ".xlsx": _extract_xlsx, ".xls": _extract_xlsx,
}

SUPPORTED_EXTENSIONS = set(_HANDLERS.keys())


def extract_text(path: str) -> str:
    """Extract plain text from a document. Raises ExtractionError on failure."""
    p = Path(path)
    handler = _HANDLERS.get(p.suffix.lower())
    if not handler:
        raise ExtractionError(f"unsupported file type: {p.suffix}")
    text = handler(p)
    text = text.strip()
    if not text:
        raise ExtractionError("no extractable text found in document")
    return text[:_MAX_CHARS]
